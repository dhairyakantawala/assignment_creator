import os
from flask import Flask, render_template, request, send_file
from langchain_openai import ChatOpenAI
from langchain.embeddings import OpenAIEmbeddings
from langchain.vectorstores import Chroma
from langchain.document_loaders import WebBaseLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.tools.retriever import create_retriever_tool
from langchain.agents import create_tool_calling_agent, AgentExecutor
from langchain import hub
from pptx import Presentation

app = Flask(__name__)

@app.route('/', methods=['GET', 'POST'])
def index():
    if request.method == 'POST':
        link = request.form['link']
        question = request.form['question']

        # Set OpenAI API Key
        os.environ["OPENAI_API_KEY"] = "sk-dqEYbbg8n1q8F7mD5sa-iBAHKM4Rh--IuFVB7UoweiT3BlbkFJhF8XCeMGnoK-bFebT9C_HBikztyFyHquuwW1X07IYA"
        os.environ["LANGCHAIN_TRACING_V2"] = "true"
        os.environ["LANGCHAIN_API_KEY"] = "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"

        # Initialize the LLM
        llm = ChatOpenAI(model="gpt-4o-mini")

        # Load the webpage content
        loader = WebBaseLoader([link])
        docs = loader.load()

        # Split the content into chunks
        text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
        split_docs = text_splitter.split_documents(docs)

        # Create embeddings and store them in Chroma
        embeddings = OpenAIEmbeddings()
        vectorstore = Chroma.from_documents(documents=split_docs, embedding=embeddings)

        # Retrieve relevant chunks based on a question
        retriever = vectorstore.as_retriever()
        retriever_tool = create_retriever_tool(
            retriever,
            "news_data",
            "Search for information about the news. For any questions about the news, you must use this tool!",
        )

        tools = [retriever_tool]
        model = ChatOpenAI(model="gpt-4o-mini", temperature=0)

        # Get the prompt to use
        prompt = hub.pull("hwchase17/openai-functions-agent")
        model_with_tools = model.bind_tools(tools)

        # Create the agent
        agent = create_tool_calling_agent(model, tools, prompt)
        agent_executor = AgentExecutor(agent=agent, tools=tools, verbose=True)

        # Generate the presentation
        ans = agent_executor.invoke({"input": f"""You are an AI assistant specialized in Python code generation. Your task is to generate Python code that creates a PowerPoint presentation using the `python-pptx` library. The code should create a presentation with slides based on titles and content that you will derive from the provided news data.

        Here are the details:
        - Use the `news_data` tool to extract relevant news headlines and summaries.
        - {question} will guide the theme or topic of the slides to be generated.
        - The resulting code should create a PowerPoint presentation with the following structure:
          - Each slide should include a title and content derived from the news data.
          - The final presentation should be saved as 'presentation.pptx'.

        Your output should be a single line of Python code that follows this structure:
        ```python
        from pptx import Presentation
        prs = Presentation()
        slides_data = [
            ("Title", "Content_point_1"),
        ]
        for title, content in slides_data:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.placeholders[0].text = title
            slide.placeholders[1].text = content
        prs.save('presentation.pptx')
        """})

        input_string = ans['output']
        start_delimiter = "```python"
        end_delimiter = "```"

        # Extracting the code
        start_index = input_string.find(start_delimiter) + len(start_delimiter)
        end_index = input_string.rfind(end_delimiter)
        code_string = input_string[start_index:end_index].strip()

        # Execute the generated code
        exec(code_string)

        return render_template('index.html', link=link, question=question, download=True)

    return render_template('index.html')

@app.route('/download')
def download_file():
    path = "presentation.pptx"
    return send_file(path, as_attachment=True)

if __name__ == '__main__':
    app.run(debug=True)
